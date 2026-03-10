from pptx import Presentation


def scan_pptx(file_path):

    issues = []

    try:

        prs = Presentation(file_path)

        # --------------------------------------------------
        # WCAG 2.4.2 – Title presence
        # --------------------------------------------------

        if not prs.core_properties.title:
            issues.append(
                {
                    "type": "pptx_missing_title",
                    "severity": "medium",
                    "location": file_path,
                }
            )

        # --------------------------------------------------
        # WCAG 1.3.1 – Slide title structure
        # --------------------------------------------------

        for i, slide in enumerate(prs.slides):

            title_found = False

            for shape in slide.shapes:

                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    title_found = True
                    break

            if not title_found:
                issues.append(
                    {
                        "type": "pptx_slide_missing_title",
                        "severity": "medium",
                        "location": f"{file_path} slide {i+1}",
                    }
                )

        # --------------------------------------------------
        # WCAG 1.1.1 – Images missing alt text
        # --------------------------------------------------

        for slide_index, slide in enumerate(prs.slides):

            for shape in slide.shapes:

                if shape.shape_type == 13:  # picture

                    alt_text = shape.element.xpath(".//p:cNvPr")

                    if alt_text:

                        descr = alt_text[0].get("descr")

                        if not descr:
                            issues.append(
                                {
                                    "type": "pptx_image_missing_alt_text",
                                    "severity": "high",
                                    "location": f"{file_path} slide {slide_index+1}",
                                }
                            )

    except Exception:

        issues.append(
            {
                "type": "pptx_scan_error",
                "severity": "low",
                "location": file_path,
            }
        )

    return issues